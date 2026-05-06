"""sandbox_it3: Seed built-in Sandbox Skills

Revision ID: sandbox_it3_builtins
Revises: sandbox_it1
Create Date: 2026-05-07

Covers IT-3 of the Sandbox Provider Integration RFC — data migration that
seeds the five globally available built-in skills:

  - word-generation       (python-docx)
  - pdf-generation        (reportlab, weasyprint)
  - presentation-generation (python-pptx)
  - data-analysis         (pandas, openpyxl, numpy)
  - charts                (matplotlib, seaborn)

All are seeded with app_id=NULL, is_builtin=TRUE, runtime='python-sandbox'.
"""
import json
from datetime import datetime

from alembic import op
import sqlalchemy as sa


# revision identifiers, used by Alembic.
revision = 'sandbox_it3_builtins'
down_revision = 'sandbox_it1'
branch_labels = None
depends_on = None

_SKILL_TABLE = 'Skill'

_BUILTIN_SKILLS = [
    {
        "name": "word-generation",
        "display_name": "Word Document Generation",
        "description": (
            "Generate and manipulate Microsoft Word (.docx) documents using python-docx. "
            "Provides document creation, formatting, tables, images, and styles."
        ),
        "content": (
            "# Word Document Generation\n\n"
            "This skill enables generation of `.docx` files inside the Python sandbox.\n\n"
            "## Usage\n\n"
            "After activating this skill, import `python-docx`:\n\n"
            "```python\n"
            "from docx import Document\n"
            "doc = Document()\n"
            "doc.add_heading('My Report', 0)\n"
            "doc.add_paragraph('Hello, world!')\n"
            "doc.save('/sandbox/output.docx')\n"
            "```\n"
        ),
        "runtime": "python-sandbox",
        "dependencies": json.dumps(["python-docx>=1.1"]),
        "allowed_tools": json.dumps(["python_repl"]),
        "bootstrap_script_path": None,
        "runtime_options": None,
        "is_builtin": True,
    },
    {
        "name": "pdf-generation",
        "display_name": "PDF Document Generation",
        "description": (
            "Generate PDF files from Python using reportlab and weasyprint. "
            "Supports text, tables, images, and HTML-to-PDF conversion."
        ),
        "content": (
            "# PDF Document Generation\n\n"
            "This skill enables PDF creation inside the Python sandbox.\n\n"
            "## Usage — reportlab\n\n"
            "```python\n"
            "from reportlab.pdfgen import canvas\n"
            "c = canvas.Canvas('/sandbox/output.pdf')\n"
            "c.drawString(100, 750, 'Hello, PDF!')\n"
            "c.save()\n"
            "```\n\n"
            "## Usage — weasyprint (HTML to PDF)\n\n"
            "```python\n"
            "import weasyprint\n"
            "weasyprint.HTML(string='<h1>Hello</h1>').write_pdf('/sandbox/output.pdf')\n"
            "```\n"
        ),
        "runtime": "python-sandbox",
        "dependencies": json.dumps(["reportlab>=4.1", "weasyprint>=62"]),
        "allowed_tools": json.dumps(["python_repl"]),
        "bootstrap_script_path": None,
        "runtime_options": None,
        "is_builtin": True,
    },
    {
        "name": "presentation-generation",
        "display_name": "Presentation Generation",
        "description": (
            "Generate Microsoft PowerPoint (.pptx) presentations using python-pptx. "
            "Supports slides, layouts, shapes, charts, and images."
        ),
        "content": (
            "# Presentation Generation\n\n"
            "This skill enables `.pptx` creation inside the Python sandbox.\n\n"
            "## Usage\n\n"
            "```python\n"
            "from pptx import Presentation\n"
            "prs = Presentation()\n"
            "slide = prs.slides.add_slide(prs.slide_layouts[0])\n"
            "slide.shapes.title.text = 'My Presentation'\n"
            "prs.save('/sandbox/output.pptx')\n"
            "```\n"
        ),
        "runtime": "python-sandbox",
        "dependencies": json.dumps(["python-pptx>=0.6.23"]),
        "allowed_tools": json.dumps(["python_repl"]),
        "bootstrap_script_path": None,
        "runtime_options": None,
        "is_builtin": True,
    },
    {
        "name": "data-analysis",
        "display_name": "Data Analysis",
        "description": (
            "Perform tabular data analysis using pandas, numpy, and openpyxl. "
            "Supports CSV, Excel reading/writing, statistical analysis, and data transformation."
        ),
        "content": (
            "# Data Analysis\n\n"
            "This skill enables data analysis workflows inside the Python sandbox.\n\n"
            "## Usage\n\n"
            "```python\n"
            "import pandas as pd\n"
            "import numpy as np\n"
            "\n"
            "df = pd.read_csv('/sandbox/data.csv')\n"
            "print(df.describe())\n"
            "df.to_excel('/sandbox/output.xlsx', index=False)\n"
            "```\n"
        ),
        "runtime": "python-sandbox",
        "dependencies": json.dumps(["pandas>=2.2", "openpyxl>=3.1", "numpy>=2.0"]),
        "allowed_tools": json.dumps(["python_repl"]),
        "bootstrap_script_path": None,
        "runtime_options": None,
        "is_builtin": True,
    },
    {
        "name": "charts",
        "display_name": "Chart Generation",
        "description": (
            "Create statistical charts and visualizations using matplotlib and seaborn. "
            "Supports bar charts, line plots, scatter plots, heat maps, and more."
        ),
        "content": (
            "# Chart Generation\n\n"
            "This skill enables chart and plot generation inside the Python sandbox.\n\n"
            "## Usage\n\n"
            "```python\n"
            "import matplotlib\n"
            "matplotlib.use('Agg')  # Non-interactive backend\n"
            "import matplotlib.pyplot as plt\n"
            "import seaborn as sns\n"
            "\n"
            "data = [1, 2, 3, 4, 5]\n"
            "plt.plot(data)\n"
            "plt.title('My Chart')\n"
            "plt.savefig('/sandbox/chart.png', bbox_inches='tight')\n"
            "plt.close()\n"
            "```\n"
        ),
        "runtime": "python-sandbox",
        "dependencies": json.dumps(["matplotlib>=3.9", "seaborn>=0.13"]),
        "allowed_tools": json.dumps(["python_repl"]),
        "bootstrap_script_path": None,
        "runtime_options": None,
        "is_builtin": True,
    },
]

_NOW = datetime(2026, 5, 7, 0, 0, 0)


def upgrade() -> None:
    skill_table = sa.table(
        _SKILL_TABLE,
        sa.column('app_id', sa.Integer),
        sa.column('name', sa.String),
        sa.column('display_name', sa.String),
        sa.column('description', sa.String),
        sa.column('content', sa.Text),
        sa.column('runtime', sa.String),
        sa.column('dependencies', sa.Text),
        sa.column('allowed_tools', sa.Text),
        sa.column('bootstrap_script_path', sa.String),
        sa.column('runtime_options', sa.Text),
        sa.column('is_builtin', sa.Boolean),
        sa.column('is_frozen', sa.Boolean),
        sa.column('create_date', sa.DateTime),
    )

    rows = [
        {
            "app_id": None,
            "name": s["name"],
            "display_name": s["display_name"],
            "description": s["description"],
            "content": s["content"],
            "runtime": s["runtime"],
            "dependencies": s["dependencies"],
            "allowed_tools": s["allowed_tools"],
            "bootstrap_script_path": s["bootstrap_script_path"],
            "runtime_options": s["runtime_options"],
            "is_builtin": True,
            "is_frozen": False,
            "create_date": _NOW,
        }
        for s in _BUILTIN_SKILLS
    ]

    op.bulk_insert(skill_table, rows)


def downgrade() -> None:
    names = [s["name"] for s in _BUILTIN_SKILLS]
    op.execute(
        sa.text(
            f"DELETE FROM \"{_SKILL_TABLE}\" "  # noqa: S608
            "WHERE app_id IS NULL AND is_builtin = TRUE "
            f"AND name IN ({', '.join(':n' + str(i) for i in range(len(names)))})"
        ).bindparams(**{f"n{i}": name for i, name in enumerate(names)})
    )
